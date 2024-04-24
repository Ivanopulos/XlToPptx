import re
import openpyxl
import pandas as pd
import pywintypes
from pptx import Presentation
import msvcrt
import time
import shutil
from zipfile import ZipFile
import os
import glob
import Nadstroika as nd
print(time.time())
print(nd.newest(".+\.xlsx"))
# 1. Прочитать Excel с помощью pandas
s_n = '20'#########################################################################################################

# Находим все Excel файлы в текущей директории
excel_files = glob.glob('*.xlsx') + glob.glob('*.xls')
# Получаем список кортежей (время последнего изменения, имя файла) для каждого файла
files_with_times = [(os.path.getmtime(file), file) for file in excel_files]
# Сортируем список по времени последнего изменения (по убыванию)
newest_file = sorted(files_with_times, key=lambda x: x[0], reverse=True)[0][1]

data_file = newest_file
pptx_name = "измененная_презентация.pptx"
prz = nd.newest(".+\.pptx")#'Заготовка_Презентация_НММО_регион_2.pptx'#####################################################################

df = pd.read_excel(newest_file, sheet_name=s_n)
df.fillna("", inplace=True)

# Проверка на наличие нужных столбцов
if 'metka' not in df.columns or 'replace_value' not in df.columns:
    raise ValueError('Не найдены столбцы "metka" или "replace_value"')

def format_value(value):
    # Проверяем, является ли значение типом float или int
    if isinstance(value, (float, int)):
        # Преобразуем значение в Decimal для повышения точности округления
        ##value = Decimal(str(value)).quantize(Decimal("1.00"))
        # Форматируем значение с разделителями разрядов, преобразуя запятую в пробел
        if value % 1 == 0:  # Если число целое
            formatted_value = f"{int(value):,}".replace(',', ' ')
        else:
            # Для чисел с плавающей точкой, заменяем точку на запятую после форматирования
            formatted_value = f"{value:,}".replace(',', ' ').replace('.', ',')

        return formatted_value
    else:
        return str(value)

replace_dict = {metka: format_value(replace_value) for metka, replace_value in zip(df['metka'], df['replace_value'])}
# def format_value(value, razrad):
#     if isinstance(value, (float, int)):
#         # Если разряд равен 0, округляем до целого числа и форматируем без десятичных знаков
#         if razrad == 0:
#             formatted_value = f"{int(round(value)):,}".replace(',', ' ')
#         # Для чисел с плавающей точкой с указанным количеством знаков после запятой
#         elif isinstance(value, float):
#             format_string = f"{{:,.{razrad}f}}"
#             formatted_value = format_string.format(value)
#             formatted_value = formatted_value.replace(',', ' ').replace('.', ',')
#         # Для целых чисел, когда разрядность не равна 0
#         else:
#             formatted_value = f"{value:,}".replace(',', ' ')
#         return formatted_value
#     else:
#         return str(value)
# replace_dict = {metka: format_value(replace_value, razrad) for metka, replace_value, razrad in zip(df['metka'], df['replace_value'], df['razrad'])}

# 2. Открыть презентацию и произвести замены
presentation = Presentation(prz)#'Заготовка Слайд_НММО.pptx')#Заготовка_Презентация_НММО_расширенная2023.pptx#################

def replace_text_in_runs(runs):
    for run in runs:
        for metka, replace_value in replace_dict.items():
            run.text = run.text.replace(metka, replace_value)

def process_shapes_for_replacement(shapes):
    for shape in shapes:
        if shape.has_text_frame:
            paragraphs = shape.text_frame.paragraphs
            for paragraph in paragraphs:
                replace_text_in_runs(paragraph.runs)

        # Обработка таблиц
        elif shape.has_table:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    paragraphs = cell.text_frame.paragraphs
                    for paragraph in paragraphs:
                        replace_text_in_runs(paragraph.runs)

for slide in presentation.slides:
    # Обработка основного текста слайдов
    process_shapes_for_replacement(slide.shapes)

    # Обработка заметок лектора
    if slide.has_notes_slide:
        paragraphs = slide.notes_slide.notes_text_frame.paragraphs
        for paragraph in paragraphs:
            replace_text_in_runs(paragraph.runs)


presentation.save(pptx_name)

import os
import zipfile
import shutil
from openpyxl import load_workbook


def modify_xml_in_zip(file_path, xml_relative_path, new_content):
    # Ensure the file_path ends with '.zip' for processing
    original_file_path = file_path
    file_path_zip = file_path if file_path.endswith('.zip') else file_path + '.zip'
    if not file_path.endswith('.zip'):
        os.rename(file_path, file_path_zip)

    # Create a temporary directory for unzipping
    temp_dir = 'temp_unzip_dir'
    if not os.path.isdir(temp_dir):
        os.mkdir(temp_dir)

    # Extract the zip file
    with zipfile.ZipFile(file_path_zip, 'r') as zip_ref:
        zip_ref.extractall(temp_dir)

    # Replace the content of the specified XML file
    xml_path = os.path.join(temp_dir, xml_relative_path)
    with open(xml_path, 'w') as xml_file:
        xml_file.write(new_content)

    # Create a new zip file (overwrite the original zip file if necessary)
    with zipfile.ZipFile(file_path_zip, 'w', zipfile.ZIP_DEFLATED) as zip_ref:
        for root, dirs, files in os.walk(temp_dir):
            for file in files:
                file_path = os.path.join(root, file)
                zip_ref.write(file_path, file_path[len(temp_dir) + 1:])

    # Clean up: remove the temporary directory
    #shutil.rmtree(temp_dir)

    # If we added '.zip' for processing, rename back to original (without '.zip')
    if original_file_path != file_path_zip:
        os.rename(file_path_zip, original_file_path)

def update_embedded_excel(pptx_name, data_file):
    temp_dir = 'temp_pptx_content'

    # 1. Разархивировать содержимое .pptx во временную папку
    with zipfile.ZipFile(pptx_name, 'r') as zip_ref:
        zip_ref.extractall(temp_dir)

    # 2. Изменить встроенный Excel документ
    data_df = pd.read_excel(data_file, sheet_name=s_n)
    for index, row in data_df.iterrows():
        if str(row['A']).startswith("ppt/"):
            file_path = os.path.join(temp_dir, row['A'])

            if os.path.exists(file_path):
                with open(file_path, 'rb') as f:
                    embedded_df = pd.read_excel(f)
##
                backup_dir = os.path.join(temp_dir, "backup")
                os.makedirs(backup_dir, exist_ok=True)
                with zipfile.ZipFile(file_path, 'r') as excel_zip:
                    print(1)
                    # Ищем файл Table1.xml в архиве
                    table1_xml_path = 'xl/tables/table1.xml'
                    if table1_xml_path in excel_zip.namelist():
                        excel_zip.extract(table1_xml_path, backup_dir)
                        table1_full_path = os.path.join(backup_dir, table1_xml_path)

                        # Выводим содержимое Table1.xml на экран
                        try:
                            with open(table1_full_path, 'r', encoding='utf-8') as table1_file:
                                nstxt = table1_file.read()
                                print(nstxt)
                        except Exception as e:
                            print(f"Ошибка при чтении файла {table1_full_path}: {e}")

                rows_count = int(row['B'])
                cols_count = int(row['C'])

                # Выберите данные, которые вы хотите добавить в embedded_df
                new_data = data_df.iloc[index + 1: index + 1 + rows_count, :cols_count]
                new_data = new_data.reset_index(drop=True)

                # Объедините новые данные с embedded_df
                embedded_df = pd.concat([new_data, embedded_df], ignore_index=True)
                embedded_df = embedded_df.iloc[:rows_count, :cols_count]

                wb = load_workbook(file_path)
                ws = wb["Лист1"]
                for row_index, row1 in embedded_df.iterrows():
                    for col_index, value in enumerate(row1):
                        ws.cell(row=row_index + 1, column=col_index + 1, value=value)
                print(wb)
                print(embedded_df)
                wb.save(file_path)
                # wb.close()
                # embedded_df.to_excel(file_path, index=False)
                #
                # xml_relative_path = "xl/tables/table1.xml"
                # modify_xml_in_zip(file_path, xml_relative_path, nstxt)



    # 3. Заархивировать содержимое временной папки
    with zipfile.ZipFile(pptx_name, 'w') as zip_ref:
        for foldername, subfolders, filenames in os.walk(temp_dir):
            for filename in filenames:
                file_path = os.path.join(foldername, filename)
                archive_path = file_path[len(temp_dir) + 1:]
                zip_ref.write(file_path, archive_path)

    # Удаляем временную папку
    shutil.rmtree(temp_dir)





update_embedded_excel(pptx_name, data_file)


import win32com.client
import time
def update_powerpoint_charts(pptx_path):
    # Создаем COM-объект для PowerPoint
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")

    # Открываем презентацию
    presentation = powerpoint.Presentations.Open(pptx_path)

    for slide in presentation.Slides:
        for shape in slide.Shapes:
            if shape.HasChart:
                # Здесь мы пытаемся "обновить" диаграмму. В вашем случае это может и не привести к результату,
                # но это максимум, что можно сделать средствами COM-объекта PowerPoint
                try:
                    shape.Chart.ChartData.Activate()
                except pywintypes.com_error as e:
                    print(e)
                    print("Не удалось активировать данные диаграммы.")
                #shape.Chart.ChartData.Activate()
                time.sleep(0.5)
                shape.Chart.ChartData.BreakLink()
                time.sleep(0.5)
                shape.Chart.Refresh()

                print(f"Обработка диаграммы на слайде {slide.SlideNumber}, название диаграммы: {shape.Name}")

    # Сохраняем и закрываем презентацию
    presentation.Save()
    presentation.Close()

    # Закрываем PowerPoint
    powerpoint.Quit()

current_directory = os.path.dirname(os.path.abspath(__file__))
pptx_path = os.path.join(current_directory, pptx_name)
print(pptx_path)
update_powerpoint_charts(pptx_path)
print(time.time())
msvcrt.getch()

